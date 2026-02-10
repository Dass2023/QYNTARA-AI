import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# --- CONSTANTS ---
# Deep Navy/Black Theme
BG_COLOR = RGBColor(1, 11, 21)       
TEXT_COLOR = RGBColor(230, 230, 230) 
ACCENT_BLUE = RGBColor(0, 162, 255)  # Tech Blue
ACCENT_GREEN = RGBColor(57, 255, 20) # Success Green
ACCENT_PURPLE = RGBColor(157, 0, 255) # Deep Creative Purple

# Update this path to the actual location
LOGO_TEXT_PATH = r"C:/Users/ArockiadassD_2j5cw54/.gemini/antigravity/brain/a616fef2-b2dc-46d8-8904-cfe20f67800b/uploaded_image_0_1766023142587.png"
LOGO_ICON_PATH = r"C:/Users/ArockiadassD_2j5cw54/.gemini/antigravity/brain/a616fef2-b2dc-46d8-8904-cfe20f67800b/uploaded_image_1766030299506.png"

def apply_slide_design(slide, is_title=False):
    """
    Applies the geometric background and logo.
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    shapes = slide.shapes
    
    # 1. Tech "Grid" Line (Top Strip)
    strip = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.1))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT_BLUE
    strip.line.fill.background() # No border
    
    # 2. Geometric Accent (Bottom Right Corner)
    triangle = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11), Inches(6), Inches(3), Inches(2))
    triangle.rotation = 45
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(10, 20, 40) # Subtle highlight
    triangle.line.fill.background()
    
    # 3. Logo Placement
    if is_title:
        # Large Centered Group for Title
        # Icon Centered
        if os.path.exists(LOGO_ICON_PATH):
           slide.shapes.add_picture(LOGO_ICON_PATH, Inches(5.9), Inches(1.5), height=Inches(1.5))
        
        # Text Logo Below Icon
        if os.path.exists(LOGO_TEXT_PATH):
           slide.shapes.add_picture(LOGO_TEXT_PATH, Inches(4.3), Inches(3.2), height=Inches(1.0))
    else:
        # Small Top-Right Icon for Content
        if os.path.exists(LOGO_ICON_PATH):
            shapes.add_picture(LOGO_ICON_PATH, Inches(12.3), Inches(0.2), height=Inches(0.6))

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 = Blank
    apply_slide_design(slide, is_title=True)
    
    shapes = slide.shapes
    
    # Title Text (Behaves like Word Art)
    title = shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.size = Pt(54)
    p.font.name = "Arial Black"
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub = shapes.add_textbox(Inches(2), Inches(5.5), Inches(9.33), Inches(1))
    tf_sub = sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = TEXT_COLOR
    p_sub.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title_str, points):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    apply_slide_design(slide, is_title=False)
    
    shapes = slide.shapes
    
    # Title
    tbox = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(1))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title_str
    p.font.size = Pt(36)
    p.font.name = "Arial"
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Accent Line under title
    line = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_GREEN
    line.line.fill.background()
    
    # Content Box
    cbox = shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf_c = cbox.text_frame
    tf_c.word_wrap = True
    
    for point in points:
        bullet = tf_c.add_paragraph()
        bullet.text = f"•  {point}"
        bullet.font.size = Pt(22)
        bullet.font.color.rgb = TEXT_COLOR
        bullet.space_after = Pt(20)

def add_architecture_visual(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, is_title=False)
    
    shapes = slide.shapes
    
    # Title
    tbox = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(1))
    tbox.text_frame.paragraphs[0].text = "Modern Hybrid Architecture"
    tbox.text_frame.paragraphs[0].font.size = Pt(36)
    tbox.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    
    # --- DIAGRAM ---
    
    # 1. MAYA Client (Left)
    maya_box = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(3), Inches(3.5))
    maya_box.fill.solid()
    maya_box.fill.fore_color.rgb = RGBColor(20, 30, 50)
    maya_box.line.color.rgb = ACCENT_BLUE
    maya_box.line.width = Pt(2)
    maya_box.text = "ARTIST\n(Maya Client)"
    maya_box.text_frame.paragraphs[0].font.bold = True
    
    # 2. CLOUD Server (Right)
    cloud_box = shapes.add_shape(MSO_SHAPE.CLOUD, Inches(8.5), Inches(2.5), Inches(3.5), Inches(3.5))
    cloud_box.fill.solid()
    cloud_box.fill.fore_color.rgb = RGBColor(40, 20, 60)
    cloud_box.line.color.rgb = ACCENT_PURPLE
    cloud_box.line.width = Pt(2)
    cloud_box.text = "QYNTARA CLOUD\n(Validator Farm)"
    
    # 3. Connector
    arrow = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.8), Inches(4.1), Inches(0.8))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
    arrow.text = "Secure HTTPS / API"
    
    # 4. Details
    info = shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1))
    info.text_frame.paragraphs[0].text = "Hybrid Integration: Instant feedback locally, massive scalability remotely."
    info.text_frame.paragraphs[0].font.size = Pt(16)
    info.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    info.text_frame.paragraphs[0].font.italic = True

def generate():
    # Widescreen Setup
    prs = Presentation() 
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slides
    add_title_slide(prs, "Qyntara AI", "The Future of 3D Asset Assurance")
    
    add_content_slide(prs, "The Challenge: Quality Bottleneck", [
        "Manual validation is slow, subjective, and prone to error.",
        "Incorrect assets break game builds and delay releases.",
        "Fixing technical debt late in production costs 10x more.",
        "Outsource consistency is hard to maintain."
    ])
    
    add_content_slide(prs, "The Solution: Qyntara AI", [
        "Automated Plugin: Runs directly in Autodesk Maya.",
        "Rule-Based Engine: 42+ deterministic checks (Geometry, Animation, Baking).",
        "AI Analysis: Deep Learning for subjective quality assessment.",
        "One-Click Auto-Fix: Resolves 100% of defined technical errors (Smart UV2 + Safe Rigging)."
    ])
    
    add_architecture_visual(prs)

    add_content_slide(prs, "Full Technology Stack", [
        "Core: Autodesk Maya 2022-2025, Python 3.9, PySide2/6.",
        "Backend: FastAPI (Rest API), Redis Queue, Docker.",
        "Geometry: OpenMaya (C++ Wrapper), NumPy, Trimesh.",
        "AI Engine: PyTorch (MeshAnomalyNet Point Cloud Model)."
    ])
    
    add_content_slide(prs, "Return on Investment (ROI)", [
        "Speed: Validation time from 2 hours -> 5 seconds.",
        "Confidence: 100% technical compliance guaranteed.",
        "Scalability: Cloud backend handles unlimited bulk processing.",
        "Workflow: Artists stay compatible with Unity/Unreal automatically."
    ])
    
    add_content_slide(prs, "Future Roadmap", [
        "Q3 2025: Multi-DCC Support (Blender, 3ds Max).",
        "Q4 2025: CI/CD Pipeline Integration.",
        "Q2 2026: Generative Geometry Repair (Auto-Retopology)."
    ])
    
    add_content_slide(prs, "Advanced AI Workflow Vision", [
        "Product Vision: Incorporate scene analysis capabilities that evaluate proportions, scale accuracy, and deviations.",
        "Precision Benchmarking: Compare assets against analytic standards and principled placement guidelines.",
        "Goal: Ensure precise and coherent object arrangement with minimal manual intervention.",
        "Scene Analysis: Verify correct proportions/scale and detect variations in real-time."
    ])

    out_path = os.path.join(os.getcwd(), "Qyntara_Management_Presentation_Pro_Custom_Logos.pptx")
    # Save
    prs.save(out_path)
    print(f"Generated Professional Deck: {out_path}")

if __name__ == "__main__":
    generate()
