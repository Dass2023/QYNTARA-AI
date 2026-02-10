import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

print("STARTING PPT GENERATION...", flush=True)

# --- CONFIG ---
WIDTH_INCHES = 13.333
HEIGHT_INCHES = 7.5

BG_COLOR = RGBColor(1, 11, 21)       # Deep Black/Blue from Logo
TEXT_COLOR = RGBColor(255, 255, 255) # White
ACCENT_1 = RGBColor(0, 206, 201)     # Cyan
ACCENT_2 = RGBColor(46, 160, 67)     # Green
ACCENT_3 = RGBColor(210, 153, 34)    # Gold

# --- IMAGE PATHS ---
BASE_DIR = os.path.join(os.getcwd(), "ppt_images")

# Map keys to partial filenames to find them dynamically
IMG_MAP = {
    "validation": "qyntara_core_validation_ui",
    "autofix": "qyntara_auto_fix_success",
    "baking": "qyntara_baking_studio",
    "uv": "qyntara_uv_toolkit",
    "ai": "qyntara_ai_alignment",
    "export": "qyntara_game_export"
}

def find_image(key):
    # Simple search in BASE_DIR
    partial = IMG_MAP.get(key)
    if not partial: return None
    
    try:
        for f in os.listdir(BASE_DIR):
            if partial in f and f.endswith(".png"):
                return os.path.join(BASE_DIR, f)
    except Exception as e:
        print(f"Warning: Could not list dir for images: {e}")
    return None

def apply_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_theme(slide)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

def add_feature_slide(prs, title, points, img_key=None, accent_color=ACCENT_1):
    # Use blank layout to control positioning
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    apply_theme(slide)
    
    # Title
    t_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
    tf = t_shape.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(40)
    p.font.color.rgb = accent_color
    p.font.bold = True
    
    # Bullets (Left Side)
    b_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    tf = b_shape.text_frame
    tf.word_wrap = True
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(14)
        p.level = 0
        if point.startswith("-"):
             p.level = 1
             p.text = point[1:].strip()
             p.font.size = Pt(18)

    # Image (Right Side)
    img_path = find_image(img_key)
    if img_path and os.path.exists(img_path):
        print(f"Embedding image: {img_path}")
        # 6.5 inches from left, 1.5 inches down, 6 inches wide
        slide.shapes.add_picture(img_path, Inches(6.8), Inches(1.8), width=Inches(6.0))

def generate():
    prs = Presentation()
    prs.slide_width = Inches(WIDTH_INCHES)
    prs.slide_height = Inches(HEIGHT_INCHES)
    
    # 1. Title
    add_title_slide(prs, "Qyntara AI: Comprehensive Overview", "End-to-End Asset Assurance Pipeline\nFeature Breakdown & Demo Plan")
    
    # 2. Intro Text
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_theme(slide)
    slide.shapes.title.text = "What Qyntara AI Covers"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "The tool is designed as an end-to-end Asset Assurance Pipeline for game development (Unreal/Unity).\n\nIt covers the entire lifecycle:\n1. Validation (Health Check)\n2. Auto-Fix (Repair)\n3. Baking (Virtual Studio)\n4. UVs (Optimization)\n5. Export (Game Ready)"
    for p in tf.paragraphs:
        p.font.color.rgb = TEXT_COLOR
        p.font.size = Pt(24)

    # 3. Core Validation
    add_feature_slide(prs, "1. Core Validation (The Health Check)", [
        "Pipeline Profiles: Auto-switch rules for Game, VR, VFX, Web.",
        "Geometry: Checks N-gons, Open Edges, Non-Manifold.",
        "Transforms: Validates Freeze, Pivot Zero, Scale.",
        "UVs: Detects Overlaps, Bounds issues.",
        "Visualizer: Viewport HUD highlights errors (Red markers).",
        "Reporting: Generates Summary HTML."
    ], "validation")

    # 4. Auto-Fix
    add_feature_slide(prs, "2. Smart Auto-Fix (One-Click)", [
        "Automated Cleanup: Deterministic repairs.",
        "Scope: Fixes History, Transforms, Unlocks Normals, Triangulates N-gons.",
        "Safety: Wrapped in Undo Chunks (Non-destructive).",
        "Result: Instantly turns a 'Red' asset to 'Green'."
    ], "autofix", accent_color=ACCENT_2)

    # 5. Baking
    add_feature_slide(prs, "3. Virtual Studio (Baking)", [
        "Virtual Rig: Procedural 3-Point Lighting (Key/Fill/Rim).",
        "AI Resolution: Suggests map size based on object scale.",
        "Production Bake: AO & Lightmaps via Arnold (GPU/CPU).",
        "Packing: Heuristic & Tetris algorithms for Lightmap UVs."
    ], "baking", accent_color=ACCENT_3)

    # 6. UV Toolkit
    add_feature_slide(prs, "4. UV Toolkit", [
        "Texel Density: Standardize to 10.24 px/cm.",
        "Smart Unwrap: AI-assisted unwrap (Angle vs Projection).",
        "Utilities: Quick Checker Map, Reset Materials.",
        "Optimization: Ensures consistent texture quality."
    ], "uv")

    # 7. Alignment / AI
    add_feature_slide(prs, "5. Alignment & AI (Experimental)", [
        "Smart Snapping: Closes gaps between modular assets.",
        "MeshAnomalyNet: Deep Learning for subjective topology analysis.",
        "Dataset Gen: Procedural training data creation.",
        "Innovation: Goes beyond simple math rules."
    ], "ai")

    # 8. Export
    add_feature_slide(prs, "6. Game Export", [
        "Engine Presets: Unreal (Z-Up) / Unity (Y-Up).",
        "Safe Export: Auto-Sanitizes before writing file.",
        "Format: Standard FBX compliant with engine importers.",
        "Workflow: 'Green Check' required before export."
    ], "export")

    # 9. Demo Script
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_theme(slide)
    slide.shapes.title.text = "Suggested Demo Workflow"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
    
    body = slide.placeholders[1]
    tf = body.text_frame
    steps = [
        "Step 1: Import 'Bad' Asset -> Show Low Health Score.",
        "Step 2: Validate -> Highlight Red Errors in Viewport.",
        "Step 3: Auto-Fix -> Watch it turn Green (100%).",
        "Step 4: Baking Tab -> 'Create Studio Rig' + 'AI Resolution'.",
        "Step 5: Export Tab -> 'Unreal Engine' -> Safe Export.",
    ]
    
    tf.clear()
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(20)

    # Save
    out_path = os.path.join(os.getcwd(), "Qyntara_Feature_Overview_With_Visuals.pptx")
    prs.save(out_path)
    print(f"Saved presentation to {out_path}", flush=True)

if __name__ == "__main__":
    try:
        generate()
    except Exception as e:
        print(f"FATAL ERROR: {e}", flush=True)
