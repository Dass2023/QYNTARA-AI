import sys
import subprocess
import os

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = subtitle_text

def add_bullet_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point

def add_architecture_diagram_slide(prs):
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "High-Level System Architecture"
    
    # Define colors
    blue = RGBColor(0, 112, 192)
    green = RGBColor(0, 176, 80)
    grey = RGBColor(128, 128, 128)
    
    # Client Side Box
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(3.0)
    height = Inches(4.0)
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.text = "Local Environment\n(Artist Workstation)"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 240, 255)
    shape.line.color.rgb = blue
    
    # Maya Node
    shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.0), Inches(2.5), Inches(0.8)).text = "Autodesk Maya\n(Client Plugin)"
    
    # Core Node
    shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(4.0), Inches(2.5), Inches(1.5)).text = "Core Logic\n- Validator\n- Auto-Fixer\n- AI Inference"

    # Arrow to Cloud
    arrow = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), Inches(4.0), Inches(1.5), Inches(0.5))
    arrow.text = "HTTPS / API"
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = grey
    
    # Cloud Side Box
    left = Inches(5.5)
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.text = "Cloud Backend\n(Docker Microservices)"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 230, 255)
    shape.line.color.rgb = RGBColor(112, 48, 160)
    
    # Components
    shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.75), Inches(2.8), Inches(2.5), Inches(0.6)).text = "FastAPI Gateway"
    shapes.add_shape(MSO_SHAPE.CAN, Inches(5.75), Inches(3.6), Inches(2.5), Inches(0.6)).text = "Redis Queue"
    shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.75), Inches(4.4), Inches(2.5), Inches(1.2)).text = "Headless Worker\n(Maya Standalone)"

def add_ai_slide(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "AI Model: MeshAnomalyNet"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "Architecture Type: PointNet-based Binary Classifier"
    
    p = tf.add_paragraph()
    p.text = "Purpose: Detects non-determinist topology errors (subjective quality)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Input: Point Cloud (N=1024 sampled vertex points)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Layers:"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "1. Convolutions (64 -> 128 -> 1024 filters) for feature extraction."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "2. Global Max Pooling (aggregates features)."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "3. Fully Connected Layers (512 -> 256 -> 1) for classification."
    p.level = 2

def add_detailed_ai_slides(prs):
    # Slide: What is an Anomaly?
    add_bullet_slide(prs, "AI Context: What is a Mesh Anomaly?", [
        "Beyond simple corruption, anomalies include:",
        "- Non-manifold geometry & Open edges",
        "- N-gons (>4 edges) & Degenerate triangles",
        "- Self-intersections & Overlapping UVs",
        "- Shadow terminator risk (Shading artifacts)",
        "- High vertex valence spikes (Rigging issues)"
    ])

    # Slide: How it Works
    add_bullet_slide(prs, "MeshAnomalyNet: How It Works", [
        "1. Representation Learning",
        "   - Converts mesh to Point Cloud or Graph (GNN)",
        "2. Neural Network Processing",
        "   - Uses PointNet++ or MeshCNN architectures",
        "   - Convolutions extract geometric features",
        "3. Detection & Output",
        "   - Classification (Good/Bad)",
        "   - Segmentation (Heatmap of defects)",
        "   - Reconstruction Error (Autoencoder pathway)"
    ])

    # Slide: Value & Applications
    add_bullet_slide(prs, "Business Value & Applications", [
        "Why use AI for Validation?",
        "- Automates Quality Control (24/7 reliability)",
        "- Reduces manual inspection time",
        "- Prevents rendering artifacts in Game Engines",
        "Target Industries:",
        "- Game Studios & VFX Houses",
        "- 3D Printing & CAD",
        "- Metaverse & UGC Pipelines"
    ])

def generate_ppt():
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, "Qyntara AI", "Next-Gen 3D Asset Validation & Auto-Fix System\nEnd-to-End Architecture Overview")
    
    # Slide 2: Executive Summary
    add_bullet_slide(prs, "Executive Summary", [
        "Hybrid System: Combines Local Maya Plugin with Cloud processing.",
        "Goal: Ensure Game/VR/VFX assets are technically flawless.",
        "Key Innovations:",
        "- 40+ Real-time Geometry Checks (N-gons, Open Edges)",
        "- Deterministic Auto-Fix Engine",
        "- Deep Learning Model (MeshAnomalyNet) for subjective quality",
        "- Dockerized Microservices for scalability"
    ])
    
    # Slide 3: Architecture Diagram (Visual)
    add_architecture_diagram_slide(prs)

    # Slide 3.5: Full Tech Stack
    add_bullet_slide(prs, "Full Technology Stack", [
        "Core: Autodesk Maya 2022-2025, Python 3.9, PySide2/6 (Qt).",
        "Backend: FastAPI (REST API), Redis (Queue), Docker containers.",
        "Geometry Engine: OpenMaya (C++ Wrapper), NumPy, Trimesh, Open3D.",
        "AI/ML: PyTorch (Deep Learning), Custom PointNet Architecture.",
        "DevOps: Git, Docker Compose, PyTest (Unit Testing)."
    ])
    
    # Slide 4: Client Side Details
    add_bullet_slide(prs, "Client Architecture (Maya Plugin)", [
        "Tech Stack: Python 3, PySide2 (Qt), Maya Python API.",
        "UI Layer: Decoupled from logic, strictly handles user signals.",
        "Core Logic: Shared library ('qyntara_ai.core') containing all validation rules.",
        "Performance: Optimized viewport overlays using OpenMaya API.",
        "Features:",
        "- Validation Dashboard",
        "- Smart Alignment Tools",
        "- Game Engine Export Presets (Unity/Unreal)"
    ])
    
    # Slide 5: Backend Details
    add_bullet_slide(prs, "Cloud/Backend Architecture", [
        "Tech Stack: FastAPI, Redis, Docker, Headless Maya.",
        "API Gateway: Handles file uploads and job validation.",
        "Message Broker: Redis queues valid jobs for processing.",
        "Worker Nodes: Scalable containers running Maya Standalone.",
        "Parity: Reuses the exact same 'Core' rules set as the client.",
        "Purpose: Allows batch processing of thousands of assets without tying up artist workstations."
    ])
    
    # Slide 6: AI
    add_ai_slide(prs)
    add_detailed_ai_slides(prs)
    
    # Slide 7: Roadmap
    add_bullet_slide(prs, "Future Roadmap", [
        "CI/CD Integration: Automated checks on Git commit.",
        "Multi-DCC Support: Adapting Core logic for Blender and 3ds Max.",
        "Generative Auto-Fix: Using AI to propose mesh topology repairs.",
        "Cloud Dashboard: Web-based viewer for validation reports."
    ])
    
    output_path = os.path.join(os.getcwd(), "Qyntara_AI_Architecture_Full.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    generate_ppt()
